"""
SiteSense — HackASU 2025 Pitch Deck
7 slides, Track 3: Economic Empowerment
Color palette: Ocean Gradient (navy #065A82, teal #1C7293, mint #02C39A)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colors ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x06, 0x5A, 0x82)
TEAL   = RGBColor(0x1C, 0x72, 0x93)
MINT   = RGBColor(0x02, 0xC3, 0x9A)
DARK   = RGBColor(0x21, 0x29, 0x5C)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF0, 0xF7, 0xFF)
RED    = RGBColor(0xE7, 0x4C, 0x3C)
YELLOW = RGBColor(0xF3, 0x9C, 0x12)
GREEN  = RGBColor(0x27, 0xAE, 0x60)
GRAY   = RGBColor(0x95, 0xA5, 0xA6)

W, H = Inches(13.3), Inches(7.5)  # LAYOUT_WIDE


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def bg(slide, color):
    """Fill slide background with solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    """Add a colored rectangle."""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def bullet_box(slide, items, x, y, w, h, size=14, color=WHITE, dot_color=None):
    """Add a bullet list."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"▸  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color


def footer(slide, text="Claude Builder Club  ·  HackASU  ·  March 20–22, 2026"):
    """Add footer bar."""
    box(slide, 0, 7.1, 13.3, 0.4, fill_color=DARK)
    txt(slide, text, 0.3, 7.15, 12.7, 0.35,
        size=9, color=GRAY, align=PP_ALIGN.CENTER)


# ── Slide 1: Title ──────────────────────────────────────────────────────────
def slide1(prs):
    s = blank_slide(prs)
    bg(s, DARK)

    # Full-width navy header band
    box(s, 0, 0, 13.3, 3.8, fill_color=NAVY)

    # Accent line
    box(s, 0, 3.8, 13.3, 0.06, fill_color=MINT)

    # Logo / icon placeholder
    txt(s, "🏗️", 0.5, 0.6, 1.5, 2.5, size=56, align=PP_ALIGN.CENTER)

    txt(s, "SiteSense", 2.0, 0.5, 10.8, 1.4,
        size=60, bold=True, color=WHITE)
    txt(s, "AI-Powered Land Feasibility in 30 Seconds",
        2.0, 1.85, 10.8, 0.9, size=24, color=MINT)
    txt(s, "HackASU 2025  ·  Track 3: Economic Empowerment & Education",
        2.0, 2.75, 10.8, 0.7, size=14, color=LIGHT)

    txt(s, '"Drop a pin, draw your lot — get a code-compliant feasibility study in 30 seconds.\n'
           'What used to cost $50,000 and 3 weeks now takes half a minute."',
        0.8, 4.1, 11.7, 1.8, size=15, italic=True, color=LIGHT,
        align=PP_ALIGN.CENTER)

    txt(s, "Mobasher Group  ·  Arizona State University",
        0, 6.1, 13.3, 0.6, size=12, color=GRAY, align=PP_ALIGN.CENTER)

    footer(s)
    return s


# ── Slide 2: The Problem ────────────────────────────────────────────────────
def slide2(prs):
    s = blank_slide(prs)
    bg(s, RGBColor(0x12, 0x1A, 0x2E))

    box(s, 0, 0, 0.12, 7.5, fill_color=RED)  # left accent bar
    txt(s, "The Problem", 0.4, 0.25, 12, 0.75, size=36, bold=True, color=WHITE)
    txt(s, "Affordable housing is in crisis — and bad land decisions are a big reason why.",
        0.4, 1.0, 12.5, 0.6, size=16, italic=True, color=LIGHT)

    # 3 pain point cards
    cards = [
        ("$10K–$50K", "Cost of a land feasibility study from a licensed PE"),
        ("3 Weeks",   "Typical time to get a preliminary site assessment"),
        ("Fly Blind", "Most buyers make $500K decisions with zero engineering insight"),
    ]
    for i, (stat, desc) in enumerate(cards):
        cx = 0.5 + i * 4.3
        box(s, cx, 1.75, 3.9, 2.5, fill_color=NAVY,
            line_color=TEAL, line_width=Pt(1))
        txt(s, stat, cx + 0.2, 1.95, 3.5, 0.9,
            size=32, bold=True, color=MINT, align=PP_ALIGN.CENTER)
        txt(s, desc, cx + 0.2, 2.85, 3.5, 1.1,
            size=13, color=LIGHT, align=PP_ALIGN.CENTER, wrap=True)

    bullet_box(s, [
        "Land buyers, architects, and developers are NOT engineers — they can't evaluate a site alone",
        "They waste months going back-and-forth with engineers on ideas that were never feasible",
        "This kills affordable housing projects before they start",
        "Our tool closes this gap — smart feasibility BEFORE the engineer meeting",
    ], 0.4, 4.4, 12.5, 2.4, size=13, color=LIGHT)

    footer(s)
    return s


# ── Slide 3: Who It Helps ───────────────────────────────────────────────────
def slide3(prs):
    s = blank_slide(prs)
    bg(s, DARK)

    box(s, 0, 0, 13.3, 1.1, fill_color=TEAL)
    txt(s, "Who It Helps", 0.5, 0.15, 12, 0.8, size=34, bold=True, color=WHITE)

    users = [
        ("🏘️", "Affordable Housing\nDeveloper",
         "Runs feasibility before spending $50K on an engineer. Builds more, faster."),
        ("🏠", "First-Time Land\nBuyer",
         "Understands what they can build before closing. No more costly surprises."),
        ("✏️", "Architect / Designer",
         "Checks site constraints in 30 seconds. Designs for what's actually buildable."),
        ("💰", "Investor /\nLand Flipper",
         "Quick due diligence on any US parcel. Better decisions, less risk."),
    ]

    for i, (icon, title, desc) in enumerate(users):
        cx = 0.4 + i * 3.15
        box(s, cx, 1.3, 2.9, 4.5, fill_color=NAVY,
            line_color=MINT, line_width=Pt(1))
        txt(s, icon, cx + 0.3, 1.45, 2.3, 0.8, size=30, align=PP_ALIGN.CENTER)
        txt(s, title, cx + 0.15, 2.2, 2.6, 0.85,
            size=13, bold=True, color=MINT, align=PP_ALIGN.CENTER)
        txt(s, desc, cx + 0.15, 3.1, 2.6, 2.4,
            size=11, color=LIGHT, align=PP_ALIGN.CENTER, wrap=True)

    txt(s, "Track 3 Equity Angle: Engineering knowledge is gatekept behind $50K consultants.\n"
           "SiteSense makes it free and instant — for everyone.",
        0.5, 6.0, 12.3, 0.9, size=13, italic=True, color=MINT, align=PP_ALIGN.CENTER)

    footer(s)
    return s


# ── Slide 4: How It Works ───────────────────────────────────────────────────
def slide4(prs):
    s = blank_slide(prs)
    bg(s, RGBColor(0x0D, 0x1B, 0x2A))

    box(s, 0, 0, 13.3, 1.1, fill_color=NAVY)
    txt(s, "How It Works — 5 Steps", 0.5, 0.15, 12, 0.8, size=34, bold=True, color=WHITE)

    steps = [
        ("1", "Drop a Pin\n& Draw", "Mapbox satellite map\nPolygon drawing tool"),
        ("2", "Auto Data\nPull",    "USGS elevation grid\nFEMA flood + USDA soil"),
        ("3", "GIS Risk\nLayers",   "Seismic · Fire · Wetlands\nWind · Snow · Landslide"),
        ("4", "Civil\nEngine",      "ACI 350-20 + ACI 360R-10\nASCE 7-22 rule engine"),
        ("5", "Plain\nEnglish",     "Claude API translates\n→ PDF report download"),
    ]

    for i, (num, title, desc) in enumerate(steps):
        cx = 0.35 + i * 2.55
        # Circle
        box(s, cx + 0.7, 1.2, 1.1, 1.1, fill_color=MINT)
        txt(s, num, cx + 0.7, 1.2, 1.1, 1.1,
            size=28, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        # Arrow between steps
        if i < 4:
            txt(s, "→", cx + 1.82, 1.55, 0.6, 0.5,
                size=16, color=GRAY, align=PP_ALIGN.CENTER)
        txt(s, title, cx + 0.1, 2.4, 2.3, 0.85,
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, desc, cx + 0.1, 3.25, 2.3, 1.2,
            size=10, color=GRAY, align=PP_ALIGN.CENTER, wrap=True)

    # Key message
    box(s, 0.5, 4.65, 12.3, 1.6, fill_color=TEAL,
        line_color=MINT, line_width=Pt(1))
    txt(s, "Address → Polygon → Analysis → Report  in under 30 seconds",
        0.7, 4.8, 11.9, 0.6, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, "All GIS data from FREE US government APIs (USGS · FEMA · USDA · USFWS) — no subscriptions needed",
        0.7, 5.35, 11.9, 0.7, size=12, color=LIGHT, align=PP_ALIGN.CENTER)

    footer(s)
    return s


# ── Slide 5: The Engineering Moat ──────────────────────────────────────────
def slide5(prs):
    s = blank_slide(prs)
    bg(s, DARK)

    box(s, 0, 0, 13.3, 1.1, fill_color=RGBColor(0xB7, 0x45, 0x0B))
    txt(s, "The Engineering Moat — Why We're Not Just ChatGPT",
        0.4, 0.12, 12.5, 0.9, size=28, bold=True, color=WHITE)

    # Left: comparison table
    box(s, 0.3, 1.2, 5.8, 5.2, fill_color=NAVY,
        line_color=TEAL, line_width=Pt(1))
    txt(s, "Others vs. SiteSense", 0.5, 1.3, 5.4, 0.5,
        size=14, bold=True, color=MINT)

    comparisons = [
        ("Generic LLM output",      "ACI 350-20 + ACI 360R-10\nASCE 7-22 validated"),
        ("No code references",       "Every rule cites chapter\n& section"),
        ("Anyone can copy",          "Proprietary rule engine\n= our IP"),
        ("Hallucinated numbers",     "USGS/FEMA/USDA data\n= defensible output"),
        ("No Arizona knowledge",     "Caliche · Washes · Expansive\nsoil · Water adequacy"),
    ]
    for i, (bad, good) in enumerate(comparisons):
        y = 1.9 + i * 0.82
        txt(s, f"✗  {bad}", 0.5, y, 2.6, 0.7, size=10, color=RED)
        txt(s, f"✓  {good}", 3.2, y, 2.7, 0.7, size=10, color=GREEN)

    # Right: code badges
    box(s, 6.6, 1.2, 6.4, 5.2, fill_color=RGBColor(0x0D, 0x1B, 0x2A),
        line_color=TEAL, line_width=Pt(1))
    txt(s, "Code Basis", 6.8, 1.3, 6.0, 0.5,
        size=14, bold=True, color=MINT)

    codes = [
        (MINT,   "ACI 350-20 / 350R-20",   "Environmental engineering concrete structures"),
        (MINT,   "ACI 360R-10",            "Slab-on-ground design and construction"),
        (YELLOW, "ASCE 7-22",              "Wind · Seismic · Flood · Snow loads"),
        (LIGHT,  "IBC 2021",               "Soils §1803 · Flood §1612"),
        (LIGHT,  "CWA §404",               "Wetlands jurisdiction"),
    ]
    for i, (color, code, desc) in enumerate(codes):
        y = 1.95 + i * 0.87
        box(s, 6.8, y, 6.0, 0.75,
            fill_color=RGBColor(0x06, 0x5A, 0x82),
            line_color=TEAL, line_width=Pt(1))
        txt(s, code, 7.0, y + 0.05, 3.5, 0.4,
            size=11, bold=True, color=color)
        txt(s, desc, 7.0, y + 0.38, 5.6, 0.35,
            size=9, color=GRAY)

    footer(s)
    return s


# ── Slide 6: Demo Screenshots (placeholder) ─────────────────────────────────
def slide6(prs):
    s = blank_slide(prs)
    bg(s, RGBColor(0x0A, 0x14, 0x22))

    box(s, 0, 0, 13.3, 1.1, fill_color=TEAL)
    txt(s, "Live Demo", 0.5, 0.15, 12, 0.8, size=34, bold=True, color=WHITE)

    # 3 demo cards
    demos = [
        ("📍 Phoenix AZ — Flat Lot",
         "Flood Zone X · Expansive soil\nConventional slab · Low risk\nSite prep: ~$85K"),
        ("🌊 Houston TX — Flood Zone AE",
         "SFHA zone · Base flood: +4 ft\nPile / elevated foundation\nSite prep: ~$210K"),
        ("⛰️ Flagstaff AZ — Slope",
         "30% slope · 40 psf snow\nDrilled caisson foundation\nSite prep: ~$175K"),
    ]
    for i, (title, stats) in enumerate(demos):
        cx = 0.4 + i * 4.3
        box(s, cx, 1.3, 4.0, 3.8, fill_color=NAVY,
            line_color=MINT, line_width=Pt(1.5))
        txt(s, title, cx + 0.15, 1.45, 3.7, 0.65,
            size=13, bold=True, color=MINT, align=PP_ALIGN.CENTER)
        txt(s, stats, cx + 0.15, 2.15, 3.7, 2.8,
            size=12, color=LIGHT, align=PP_ALIGN.CENTER, wrap=True)

    # Feature list
    features = [
        "🗺️  Satellite map + polygon draw",
        "📊  Elevation cross-section chart",
        "🔴🟡🟢  Traffic-light risk cards (flood, seismic, fire, soil, wetlands)",
        "🏗️  Cut & fill bar chart (cubic yards)",
        "💰  Cost now + 10-year projection table",
        "📄  One-click PDF report (Claude plain-English translation)",
    ]
    bullet_box(s, features, 0.4, 5.3, 12.5, 1.8, size=11, color=LIGHT)

    footer(s)
    return s


# ── Slide 7: Why We Win ─────────────────────────────────────────────────────
def slide7(prs):
    s = blank_slide(prs)
    bg(s, NAVY)

    # Big closing statement
    box(s, 0, 0, 13.3, 2.8, fill_color=DARK)
    txt(s, "Why We Win", 0.5, 0.2, 12, 0.8, size=36, bold=True, color=MINT)
    txt(s, '"You bought a piece of land. Now what? Most developers waste months and thousands of dollars\n'
           'going back-and-forth with engineers on ideas that were never feasible.\n'
           'Our tool changes that — in 30 seconds."',
        0.5, 0.95, 12.3, 1.7, size=14, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

    # 3 win pillars
    pillars = [
        ("🎯", "Real Problem",
         "$50K feasibility → free in 30 sec\nMillions of land deals per year\nAffordable housing crisis"),
        ("🔬", "Technical Moat",
         "ACI 350-20 + ACI 360R-10\nASCE 7-22 rule engine\n10 GIS data layers\nArizona-specific rules"),
        ("🚀", "Beyond Hackathon",
         "Expand to full USA zoning\nEngineer marketplace\nSell API to Zillow / Redfin\nThailand + SE Asia next"),
    ]
    for i, (icon, title, desc) in enumerate(pillars):
        cx = 0.5 + i * 4.25
        box(s, cx, 3.0, 3.9, 3.7, fill_color=TEAL,
            line_color=MINT, line_width=Pt(1.5))
        txt(s, icon, cx + 0.3, 3.1, 3.3, 0.8, size=32, align=PP_ALIGN.CENTER)
        txt(s, title, cx + 0.15, 3.9, 3.6, 0.6,
            size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, desc, cx + 0.15, 4.5, 3.6, 2.0,
            size=11, color=LIGHT, align=PP_ALIGN.CENTER, wrap=True)

    txt(s, "Mobasher Group · ASU Civil Engineering · HackASU 2025",
        0, 6.85, 13.3, 0.5, size=11, color=MINT, align=PP_ALIGN.CENTER)

    footer(s)
    return s


# ── Main ─────────────────────────────────────────────────────────────────────
def build_deck():
    prs = new_prs()
    prs.core_properties.title = "SiteSense — AI Land Feasibility Tool"
    prs.core_properties.author = "Mobasher Group, ASU"

    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    slide5(prs)
    slide6(prs)
    slide7(prs)

    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "..", "HackASU_SiteSense_Pitch_Deck.pptx")
    prs.save(out_path)
    print(f"Saved: {os.path.abspath(out_path)}")
    return out_path


if __name__ == "__main__":
    build_deck()
